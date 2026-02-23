import plotly.io as pio
import streamlit as st
import pandas as pd
import plotly.express as px
import numpy as np

# Detect current theme
theme_type = st.context.theme.type if hasattr(st.context.theme, 'type') else "light"  # fallback

# Define schemes (your colors)
if theme_type == "dark":
    PRIMARY = "#A3CDED"          # brighter for dark visibility
    ACCENT = "#42BA97"           # teal keeps good contrast
    SAVINGS_GREEN = "#66bb6a"    # lighter green
    LOSS_RED = "#ff7474"         # visible red
    NEUTRAL = "#1E293B"          # lighter gray for text/averages
    PLOTLY_TEMPLATE = "plotly_dark"
else:  # light or unknown/system-light
    PRIMARY = "#335B74"
    ACCENT = "#42BA97"
    SAVINGS_GREEN = "#3E8853"
    LOSS_RED = "#FF7474"
    NEUTRAL = "#808080"
    PLOTLY_TEMPLATE = "plotly_white"

# Optional: small CSS injection for anything config.toml misses (e.g., specific titles)
st.markdown(f"""
    <style>
        h1, h2, h3 {{ color: {PRIMARY} !important; }}
        /* Add more if needed, but keep minimal */
    </style>
""", unsafe_allow_html=True)

# ── PPTX Export Function (fixed - no longer requires savings_to_cms_m) ────────
def generate_pptx_report(aco_data, df, track_avg):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 wide layout
    prs.slide_height = Inches(7.5)

    # ── Colors ───────────────────────────────────────────────────────────────
    from pptx.dml.color import RGBColor
    def hex_to_rgb(h): return tuple(int(h.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
    primary_rgb = hex_to_rgb(PRIMARY) if PRIMARY.startswith('#') else (51, 91, 116)
    text_rgb    = (40, 40, 40)
    gray_rgb    = (120, 120, 120)
    light_bg_rgb = (250, 250, 252)

    # ── Helpers ───────────────────────────────────────────────────────────────
    def set_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*light_bg_rgb)

    def set_footer(slide):
        left, top, w, h = Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3)
        txBox = slide.shapes.add_textbox(left, top, w, h)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"PY 2023 MSSP Dashboard | Data: CMS PUF | Generated {pd.Timestamp.now().strftime('%Y-%m-%d')}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*gray_rgb)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        set_background(slide)
        title = slide.shapes.title
        title.text = "Medicare Shared Savings Program – PY 2023 Dashboard Report"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_rgb)
        subtitle = slide.placeholders[1]
        subtitle.text = (
            f"Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M CST')}\n"
            "Chicago, Illinois | @ASmithson1987\n"
            f"Selected ACO: {aco_data.get('ACO_Name', 'N/A')} ({aco_data.get('ACO_ID', 'N/A')})"
        )
        set_footer(slide)

    def add_section_title(title):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_background(slide)
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_rgb)
        set_footer(slide)
        return slide

    def add_bullet_slide(title, bullets):
        slide = add_section_title(title)
        tf = slide.placeholders[1].text_frame
        for line in bullets:
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 1 if line.startswith("- ") else 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*text_rgb)

    def add_table_slide(title, data_dict):
        slide = add_section_title(title)
        rows = len(next(iter(data_dict.values()))) + 1
        cols = len(data_dict)
        left, top, width, height = Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.28 * rows)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Headers
        for c, hdr in enumerate(data_dict):
            cell = table.cell(0, c)
            cell.text = hdr
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*primary_rgb)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(13)

        # Data + alternating rows
        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                if r == 0: continue
                cell.text = str(list(data_dict.values())[c][r-1])
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(*text_rgb)
                if r % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

        set_footer(slide)

    def add_metrics_slide(title, metrics):
        slide = add_section_title(title)
        left, top, card_w, card_h, spacing = Inches(0.7), Inches(1.4), Inches(2.3), Inches(1.6), Inches(0.3)
        for i, (label, value) in enumerate(metrics):
            x = left + i * (card_w + spacing)
            shape = slide.shapes.add_textbox(x, top, card_w, card_h)
            tf = shape.text_frame
            tf.clear()
            line = shape.line
            line.color.rgb = RGBColor(*primary_rgb)
            line.width = Pt(1.2)
            p1 = tf.add_paragraph(); p1.text = value; p1.font.bold = True; p1.font.size = Pt(30); p1.font.color.rgb = RGBColor(*primary_rgb); p1.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(13); p2.font.color.rgb = RGBColor(*text_rgb); p2.alignment = PP_ALIGN.CENTER
        set_footer(slide)

    # ── Slides ────────────────────────────────────────────────────────────────
    add_title_slide()

    add_bullet_slide("PY 2023 Program Highlights", [
        "- Shift to sliding-scale quality scoring (partial savings + health equity bonus up to 10 points)",
        "- Expanded primary care service codes for prospective assignment",
        "- Continued EUC flexibilities (all ACOs DisAffQual=1 due to COVID-19 PHE)",
        "- Record net savings to Medicare: $2.1B (largest in program history)",
        "- Quality improvements across many measures vs PY 2022"
    ])

    changes_data = {
        "Aspect": ["Quality Scoring", "EUC/DisAffQual", "Assignment Codes", "Risk/County Data", "Savings/Losses Calc"],
        "PY 2023 Handling": [
            "Sliding scale + health equity bonus possible",
            "All ACOs flagged =1 (COVID-19 PHE)",
            "Expanded primary care codes",
            "Suppression rules unchanged (1-10 suppressed)",
            "Sliding scale sharing; no MSR waiver changes"
        ],
        "Impact on Data": [
            "QualScore reflects MIPS + bonus; affects EarnSaveLoss",
            "Applies to quality/financial adjustments",
            "May increase N_AB in some ACOs",
            "Weighted risk scores use valid PY data only",
            "GenSaveLoss/EarnSaveLoss quality-adjusted"
        ],
        "Source": [
            "CY 2023 PFS Final Rule",
            "Data Dictionary (DisAffQual)",
            "Final Rule assignment updates",
            "County-Level FFS Methodology PUF",
            "Data Dictionary + Final Rule"
        ]
    }
    add_table_slide("Key Technical/Data Changes in PY 2023 PUF", changes_data)

    add_bullet_slide("Data Sources & Methodology", [
        "- Primary: PY 2023 ACO Results PUF (CMS)",
        "- All ACOs DisAffQual=1 (COVID-19 PHE)",
        "- Small cell suppression: 1–10 assignable beneficiaries suppressed",
        "- Regional FFS & risk scores per County-Level PUF",
        "- Regulations: 42 CFR Part 425",
        "- Generated from Streamlit MSSP Dashboard"
    ])

    total_generated_m = df['GenSaveLoss'].sum(skipna=True) / 1_000_000
    total_earned_m = df['EarnSaveLoss'].sum(skipna=True) / 1_000_000
    savings_to_cms_m = (df['GenSaveLoss'].sum(skipna=True) - df['EarnSaveLoss'].sum(skipna=True)) / 1_000_000
    metrics = [
        ("Total ACOs", f"{len(df):,}"),
        ("Assigned Beneficiaries", fmt_comma(df['N_AB'].sum())),
        ("Generated Savings", f"${total_generated_m:,.2f}M"),
        ("Earned by ACOs", f"${total_earned_m:,.2f}M"),
        ("Net to CMS", f"${savings_to_cms_m:,.2f}M")
    ]
    add_metrics_slide("Program-Wide Overview (PY 2023)", metrics)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()

# ── Load & Prepare Data ──────────────────────────────────────────────────────
@st.cache_data
def load_data():
    df = pd.read_csv("PY 2023 ACO Results PUF.csv", low_memory=False)
    df = df.replace({"*": np.nan, "": np.nan, "-": np.nan})

    # Convert quality & savings to numeric
    for col in ['QualScore', 'Sav_rate']:
        df[col] = pd.to_numeric(
            df[col].astype(str).str.replace("%", "", regex=False).str.strip(),
            errors="coerce"
        )

    # Numeric columns (expanded to include all demographic and beneficiary columns)
    numeric_cols = [
        "N_AB", "BnchmkMinExp", "GenSaveLoss", "EarnSaveLoss",
        "ABtotBnchmk", "ABtotExp", "Per_Capita_Exp_TOTAL_PY"
    ] + [col for col in df.columns if any(k in col for k in [
        "CapAnn_", "ADM_", "SNF_LOS", "CMS_HCC_RiskScore_", "CAHPS_", "Measure_", "QualityID_",
        "N_AB_Year_ESRD_PY", "N_AB_Year_DIS_PY", "N_AB_Year_AGED_Dual_PY", "N_AB_Year_AGED_NonDual_PY",
        "N_Ben_VA_Only", "N_Ben_CBA_Only", "N_Ben_CBA_and_VA",
        "N_Ben_Female", "N_Ben_Male",
        "N_Ben_Age_0_64", "N_Ben_Age_65_74", "N_Ben_Age_75_84", "N_Ben_Age_85plus",
        "N_Ben_Race_White", "N_Ben_Race_Black", "N_Ben_Race_Asian", "N_Ben_Race_Hisp",
        "N_Ben_Race_Native", "N_Ben_Race_Other", "N_Ben_Race_Unknown"
    ])]
    for col in numeric_cols:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    # Derived
    df["Per_Beneficiary_Savings"] = df["BnchmkMinExp"] / df["N_AB"].replace(0, np.nan)
    df["PMPM_Savings"] = df["Per_Beneficiary_Savings"] / 12
    df["HighSav_LowQual"] = (df["Sav_rate"] > 5) & (df["QualScore"] < 80)

    # Proxy total per capita utilization
    capann_major = [col for col in df.columns if col.startswith("CapAnn_")]
    df["Total_CapAnn_Util"] = df[capann_major].sum(axis=1, skipna=True)

    # PY weighted risk score (skips suppressed categories)
    def calculate_weighted_risk_py(row):
        risk_scores = {
            "ESRD": pd.to_numeric(row.get("CMS_HCC_RiskScore_ESRD_PY"), errors='coerce'),
            "Disabled": pd.to_numeric(row.get("CMS_HCC_RiskScore_DIS_PY"), errors='coerce'),
            "Aged Dual": pd.to_numeric(row.get("CMS_HCC_RiskScore_AGDU_PY"), errors='coerce'),
            "Aged Non-Dual": pd.to_numeric(row.get("CMS_HCC_RiskScore_AGND_PY"), errors='coerce')
        }

        person_years = {
            "ESRD": pd.to_numeric(row.get("N_AB_Year_ESRD_PY"), errors='coerce'),
            "Disabled": pd.to_numeric(row.get("N_AB_Year_DIS_PY"), errors='coerce'),
            "Aged Dual": pd.to_numeric(row.get("N_AB_Year_AGED_Dual_PY"), errors='coerce'),
            "Aged Non-Dual": pd.to_numeric(row.get("N_AB_Year_AGED_NonDual_PY"), errors='coerce')
        }

        valid_terms = []
        total_weight = 0.0

        for cat in risk_scores:
            risk = risk_scores[cat]
            py = person_years[cat]
            if pd.notna(risk) and pd.notna(py) and py > 0:
                valid_terms.append(risk * py)
                total_weight += py

        if total_weight > 0:
            return sum(valid_terms) / total_weight
        else:
            return np.nan  # No valid data → NaN

    df["weighted_risk_py"] = df.apply(calculate_weighted_risk_py, axis=1)

    # BY3 weighted risk score (same logic as PY)
    def calculate_weighted_risk_by3(row):
        risk_scores = {
            "ESRD": pd.to_numeric(row.get("CMS_HCC_RiskScore_ESRD_BY3"), errors='coerce'),
            "Disabled": pd.to_numeric(row.get("CMS_HCC_RiskScore_DIS_BY3"), errors='coerce'),
            "Aged Dual": pd.to_numeric(row.get("CMS_HCC_RiskScore_AGDU_BY3"), errors='coerce'),
            "Aged Non-Dual": pd.to_numeric(row.get("CMS_HCC_RiskScore_AGND_BY3"), errors='coerce')
        }

        person_years = {
            "ESRD": pd.to_numeric(row.get("N_AB_Year_ESRD_BY3"), errors='coerce'),
            "Disabled": pd.to_numeric(row.get("N_AB_Year_DIS_BY3"), errors='coerce'),
            "Aged Dual": pd.to_numeric(row.get("N_AB_Year_AGED_Dual_BY3"), errors='coerce'),
            "Aged Non-Dual": pd.to_numeric(row.get("N_AB_Year_AGED_NonDual_BY3"), errors='coerce')
        }

        valid_terms = []
        total_weight = 0.0

        for cat in risk_scores:
            risk = risk_scores[cat]
            py = person_years[cat]
            if pd.notna(risk) and pd.notna(py) and py > 0:
                valid_terms.append(risk * py)
                total_weight += py

        if total_weight > 0:
            return sum(valid_terms) / total_weight
        else:
            return np.nan  # No valid data → NaN

    df["weighted_risk_by3"] = df.apply(calculate_weighted_risk_by3, axis=1)

    return df

df = load_data()

# ── Format helpers ───────────────────────────────────────────────────────────
def fmt_dollars(x, decimals=0):
    if pd.isna(x): return "-"
    try:
        return f"${float(x):,.{decimals}f}"
    except:
        return str(x)

def fmt_pct(x):
    if pd.isna(x): return "-"
    try:
        return f"{float(x):.2f}%"
    except:
        return str(x) if str(x).strip() else "-"

def fmt_comma(x):
    if pd.isna(x): return "-"
    try:
        return f"{int(round(x)):,}"
    except:
        return str(x)

def fmt_date(x):
    if pd.isna(x): return "-"
    try:
        return pd.to_datetime(x).strftime("%Y-%m-%d")
    except:
        return str(x)

# ── Final robust fixed banner ────────────────────────────────────────────────
st.set_page_config(page_title="MSSP 2023 Dashboard", layout="wide")

# ── Fixed banner CSS ─────────────────────────────────────────────────────────
st.markdown(f"""
    <style>
        /* Fixed banner at absolute top */
        .fixed-banner {{
            position: fixed !important;
            top: 0 !important;
            left: 0 !important;
            right: 0 !important;
            width: 100% !important;
            height: 60px !important;
            background: {'#f8f9fa' if theme_type != "dark" else '#0e1117'} !important;
            border-bottom: 1px solid {'#ddd' if theme_type != "dark" else '#333'} !important;
            box-shadow: 0 2px 6px rgba(0,0,0,0.1) !important;
            z-index: 1000 !important;
            padding: 0 20px !important;
            display: flex !important;
            align-items: center !important;
            justify-content: flex-start !important;
        }}
        
        .fixed-banner .title {{
            font-size: 1.4rem !important;
            color: {PRIMARY} !important;
            margin: 0 !important;
            padding: 0 !important;
            flex: 1 !important;
            text-align: center !important;
        }}
        
        /* Push app content below banner */
        [data-testid="stAppViewContainer"] {{
            padding-top: 70px !important;
        }}
        
        /* Push sidebar below banner */
        [data-testid="stSidebar"] {{
            top: 70px !important;
            position: absolute !important;
            height: calc(100vh - 70px) !important;
        }}
    </style>
""", unsafe_allow_html=True)

# Fixed banner markup
st.markdown("""
    <div class="fixed-banner">
        <div class="title">Medicare Shared Savings Program – PY 2023 Dashboard</div>
    </div>
""", unsafe_allow_html=True)

tab_selection = st.sidebar.radio(
    "Navigate",
    ["Program Changes", "Overview", "Single ACO View"],
    index=0,
    key="main_navigation"
)

# Main content area – switch based on selection
if tab_selection == "Overview":
    st.subheader("Program-Wide Totals")
    cols = st.columns(6)
    cols[0].metric("Total ACOs", f"{len(df):,}")
    cols[1].metric("Total Assigned Beneficiaries", f"{df['N_AB'].sum():,}")
    total_generated_m = df['GenSaveLoss'].sum(skipna=True) / 1_000_000
    cols[2].metric("Total Generated Savings ($M)", f"${total_generated_m:,.2f}M")
    total_earned_m = df['EarnSaveLoss'].sum(skipna=True) / 1_000_000
    cols[3].metric("Total Earned by ACOs ($M)", f"${total_earned_m:,.2f}M")
    savings_to_cms_m = (df['GenSaveLoss'].sum(skipna=True) - df['EarnSaveLoss'].sum(skipna=True)) / 1_000_000
    cols[4].metric("Net Savings to CMS ($M)", f"${savings_to_cms_m:,.2f}M",
                   delta_color="normal" if savings_to_cms_m >= 0 else "inverse")
    valid = df[df['Sav_rate'].notna()]
    w_avg_sav = np.average(valid['Sav_rate'], weights=valid['N_AB']) if len(valid) > 0 else np.nan
    cols[5].metric("Weighted Avg Savings Rate", f"{w_avg_sav:.2f}%" if not np.isnan(w_avg_sav) else "N/A")
    st.markdown("---")
    view_mode = st.radio("View aggregates as:", ["Totals", "Per Beneficiary", "PMPM"], horizontal=True)
    agg = df.groupby("Current_Track").agg(
        ACO_Count=("ACO_ID", "count"),
        Beneficiaries=("N_AB_Year_PY", "sum"),
        Benchmark=("ABtotBnchmk", "sum"),
        Expenditures=("ABtotExp", "sum"),
        Raw_Savings=("BnchmkMinExp", "sum"),
        Earned=("EarnSaveLoss", "sum"),
    ).reset_index()
    agg = agg.rename(columns={"Beneficiaries": "Assigned Beneficiaries"})
    agg["Weighted_Savings_Rate"] = pd.to_numeric(
        df.groupby("Current_Track").apply(
            lambda g: np.average(g['Sav_rate'], weights=g['N_AB']) if g['Sav_rate'].notna().any() else np.nan
        ).values,
        errors="coerce"
    )
    # Always calculate per-unit columns (even if not displayed in "Totals" mode)
    agg["Per_Ben_Benchmark"] = agg["Benchmark"] / agg["Assigned Beneficiaries"]
    agg["Per_Ben_Expenditures"] = agg["Expenditures"] / agg["Assigned Beneficiaries"]
    agg["Per_Ben_Raw_Savings"] = agg["Raw_Savings"] / agg["Assigned Beneficiaries"]
    agg["Per_Ben_Earned"] = agg["Earned"] / agg["Assigned Beneficiaries"]
    agg["PMPM_Benchmark"] = agg["Per_Ben_Benchmark"] / 12
    agg["PMPM_Expenditures"] = agg["Per_Ben_Expenditures"] / 12
    agg["PMPM_Raw_Savings"] = agg["Per_Ben_Raw_Savings"] / 12
    agg["PMPM_Earned"] = agg["Per_Ben_Earned"] / 12
    # Choose columns based on view_mode, but always include the core ones
    core_cols = ["Current_Track", "ACO_Count", "Assigned Beneficiaries"]
    if view_mode == "Per Beneficiary":
        disp_cols = core_cols + ["Per_Ben_Benchmark", "Per_Ben_Expenditures", "Per_Ben_Raw_Savings", "Per_Ben_Earned", "Weighted_Savings_Rate"]
    elif view_mode == "PMPM":
        disp_cols = core_cols + ["PMPM_Benchmark", "PMPM_Expenditures", "PMPM_Raw_Savings", "PMPM_Earned", "Weighted_Savings_Rate"]
    else:
        disp_cols = core_cols + ["Benchmark", "Expenditures", "Raw_Savings", "Earned", "Weighted_Savings_Rate"]
    disp = agg[disp_cols].copy()
    # Formatting loop – now handles the new per-unit columns
    for col in disp.columns:
        if col == "Current_Track":
            continue
        if col in ["ACO_Count", "Assigned Beneficiaries"]:
            disp[col] = disp[col].apply(fmt_comma)
        elif col == "Weighted_Savings_Rate":
            disp[col] = disp[col].apply(fmt_pct)
        elif "PMPM" in col:
            disp[col] = disp[col].apply(lambda x: fmt_dollars(x, decimals=2))
        elif "Per_Ben" in col:
            disp[col] = disp[col].apply(lambda x: fmt_dollars(x, decimals=0))
        elif any(k in col for k in ["Benchmark", "Expenditures", "Raw_Savings", "Earned"]):
            disp[col] = disp[col].apply(lambda x: fmt_dollars(x, decimals=0))
    # Rename for display clarity (add per-unit labels)
    rename_dict = {
        "Per_Ben_Benchmark": "Benchmark (Per Beneficiary)",
        "Per_Ben_Expenditures": "Expenditures (Per Beneficiary)",
        "Per_Ben_Raw_Savings": "Raw Savings (Per Beneficiary)",
        "PMPM_Benchmark": "Benchmark (PMPM)",
        "PMPM_Expenditures": "Expenditures (PMPM)",
        "PMPM_Raw_Savings": "Raw Savings (PMPM)",
    }
    disp = disp.rename(columns=rename_dict)
    st.dataframe(disp, use_container_width=True, hide_index=True)
    st.markdown("---")
    st.subheader("Visual Breakdown by Track")
    colL, colR = st.columns(2)
    with colL:
        fig1 = px.bar(agg.melt(id_vars="Current_Track", value_vars=["Benchmark", "Expenditures"]),
                      x="Current_Track", y="value", color="variable", barmode="group",
                      title="Benchmark vs Expenditures", color_discrete_sequence=["#4A8AB8", "#FF7474"])
        fig1.update_layout(template=PLOTLY_TEMPLATE, yaxis_tickformat="$,.0f")
        st.plotly_chart(fig1, use_container_width=True)
    with colR:
        fig2 = px.bar(agg.melt(id_vars="Current_Track", value_vars=["Raw_Savings", "Earned"]),
                      x="Current_Track", y="value", color="variable", barmode="group",
                      title="Raw Savings vs Earned by ACOs", color_discrete_sequence=["#42BA97", "#B381D9"])
        fig2.update_layout(template=PLOTLY_TEMPLATE, yaxis_tickformat="$,.0f")
        st.plotly_chart(fig2, use_container_width=True)
    colP1, colP2 = st.columns(2)
    # Define pie colors based on current theme (reuses your variables)
    pie_colors = [PRIMARY, ACCENT, NEUTRAL, "#636EFA", "#EF553B"] # fallback + theme colors
    with colP1:
        fig_aco_pie = px.pie(
            agg,
            values="ACO_Count",
            names="Current_Track",
            title="ACOs by Track",
            hole=0.4,
            color_discrete_sequence=pie_colors # ← custom colors
        )
        fig_aco_pie.update_layout(
            template=PLOTLY_TEMPLATE,
            showlegend=True,
            legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5)
        )
        st.plotly_chart(fig_aco_pie, use_container_width=True)
    with colP2:
        fig_benef_pie = px.pie(
            agg,
            values="Assigned Beneficiaries",
            names="Current_Track",
            title="Beneficiaries by Track",
            hole=0.4,
            color_discrete_sequence=pie_colors # ← same custom colors
        )
        fig_benef_pie.update_layout(
            template=PLOTLY_TEMPLATE,
            showlegend=True,
            legend=dict(orientation="h", yanchor="bottom", y=-0.2, xanchor="center", x=0.5)
        )
        st.plotly_chart(fig_benef_pie, use_container_width=True)
    # ── Moved from Charts tab: Quality & Performance Visuals ──────────────────
    st.markdown("---")
    st.subheader("Quality & Performance Visuals (PY 2023)")
    # 1. Quality Score vs Earned Savings Scatter
    st.markdown("### 1. Quality Score vs Earned Shared Savings/Losses")
    track_options = ["All"] + sorted(df["Current_Track"].dropna().unique())
    selected_track = st.selectbox("Filter by Track (Scatter)", track_options, index=0, key="track_scatter_overview")
    plot_df = df.copy()
    if selected_track != "All":
        plot_df = plot_df[plot_df["Current_Track"] == selected_track]
    plot_df = plot_df[plot_df["QualScore"].notna() & plot_df["EarnSaveLoss"].notna()]
    if len(plot_df) == 0:
        st.warning("No ACOs with valid data match the filter.")
    else:
        plot_df["Highlight"] = np.where(
            (plot_df["Sav_rate"] > 5) & (plot_df["QualScore"] < 80),
            "High Savings + Low Quality (Outreach Candidate)",
            "Other"
        )
        fig_earned = px.scatter(
            plot_df,
            x="QualScore",
            y="EarnSaveLoss",
            size="N_AB",
            color="Highlight",
            color_discrete_map={
                "High Savings + Low Quality (Outreach Candidate)": "#FF7474",
                "Other": "#2683C6"
            },
            hover_name="ACO_Name",
            hover_data={
                "EarnSaveLoss": ":$,.0f",
                "Sav_rate": ":.2f%",
                "QualScore": ":.1f",
                "N_AB": ":,",
                "Current_Track": True
            },
            title="Quality Score vs Earned Shared Savings/Losses<br>Red = High Savings Rate (>5%) + Low Quality (<80%)",
            labels={
                "QualScore": "Quality Score",
                "EarnSaveLoss": "Earned Shared Savings/Losses ($)"
            },
            trendline="ols",
            trendline_scope="overall",
            trendline_color_override="#A5D8FF"
        )
        fig_earned.update_traces(marker=dict(opacity=0.75, line=dict(width=0.5)))
        fig_earned.update_layout(template=PLOTLY_TEMPLATE, yaxis_tickformat="$,.0f", height=600)
        st.plotly_chart(fig_earned, use_container_width=True)
    # 2. Box Plot: Quality Score by Track
    st.markdown("### 2. Quality Score Distribution by Track")
    track_order = ["A", "B", "C", "D", "E", "EN"]
    fig_box = px.box(
        df[df["QualScore"].notna()],
        x="Current_Track",
        y="QualScore",
        color="Current_Track",
        points="outliers",
        title="Distribution of Quality Scores by Track",
        labels={"Current_Track": "Track", "QualScore": "Quality Score"},
        category_orders={"Current_Track": track_order}
    )
    fig_box.update_layout(
        template=PLOTLY_TEMPLATE,
        height=500,
        showlegend=False,
        xaxis={'categoryorder':'array', 'categoryarray': track_order}
    )
    st.plotly_chart(fig_box, use_container_width=True)
    # 3. CAHPS Domain Averages by Track
    st.markdown("### 3. CAHPS Domain Averages by Track")
    cahps_domains = {
        "Getting Timely Care": ["CAHPS_1", "CAHPS_2"],
        "Doctor Communication": ["CAHPS_3", "CAHPS_4", "CAHPS_5"],
        "Patient Rating of Provider": ["CAHPS_6"],
        "Access to Specialists": ["CAHPS_7"],
        "Health Promotion & Education": ["CAHPS_11"],
        "Shared Decision Making": ["CAHPS_9"],
        "Care Coordination": ["CAHPS_8"]
    }
    cahps_cols = [col for col in df.columns if col.startswith("CAHPS_")]
    if cahps_cols:
        cahps_melt = df.melt(
            id_vars=["Current_Track"],
            value_vars=cahps_cols,
            var_name="Measure",
            value_name="Score"
        )
        cahps_melt["Score"] = pd.to_numeric(cahps_melt["Score"], errors="coerce")
        def get_domain(measure):
            for domain, measures in cahps_domains.items():
                if measure in measures:
                    return domain
            return "Other"
        cahps_melt["Domain"] = cahps_melt["Measure"].apply(get_domain)
        cahps_melt = cahps_melt[cahps_melt["Domain"] != "Other"]
        domain_avg = cahps_melt.groupby(["Current_Track", "Domain"])["Score"].mean().reset_index()
        fig_cahps = px.bar(
            domain_avg,
            x="Domain",
            y="Score",
            color="Current_Track",
            barmode="group",
            title="Average CAHPS Domain Scores by Track",
            labels={"Score": "Average Score (%)", "Domain": "CAHPS Domain"}
        )
        fig_cahps.update_layout(
            template=PLOTLY_TEMPLATE,
            yaxis_range=[0, 100],
            xaxis_tickangle=-45,
            height=600
        )
        st.plotly_chart(fig_cahps, use_container_width=True)
    else:
        st.info("CAHPS measures not available or all suppressed.")

elif tab_selection == "Single ACO View":
    st.subheader("Single ACO View")
    # ACO selector
    aco_options = sorted(df["ACO_Name"].unique())
    selected_aco = st.selectbox("Select ACO", aco_options, index=0)
    aco_data = df[df["ACO_Name"] == selected_aco].iloc[0]
    if aco_data.empty:
        st.error("No data found for selected ACO.")
    else:
        # Force beneficiary and demographic columns to numeric
        beneficiary_cols = [col for col in df.columns if any(k in col for k in ["N_AB", "N_Ben_", "N_AB_Year_", "N_Ben_VA_Only", "N_Ben_CBA_Only", "N_Ben_CBA_and_VA"])]
        for col in beneficiary_cols:
            if col in aco_data:
                aco_data[col] = pd.to_numeric(aco_data[col], errors="coerce")
        # Calculate track average early
        track = aco_data["Current_Track"]
        track_avg = df[df["Current_Track"] == track].mean(numeric_only=True)
        # 1. Header & Configuration
        st.markdown(f"### {aco_data['ACO_Name']} ({aco_data['ACO_ID']})")
        cols_row1 = st.columns(5)
        cols_row1[0].metric("Track", aco_data["Current_Track"])
        cols_row1[1].metric("Agreement Type", aco_data["Agree_Type"])
        cols_row1[2].metric("Current Start Date", fmt_date(aco_data["Current_Start_Date"]))
        cols_row1[3].metric("Assignment Method", aco_data["Assign_Type"])
        cols_row1[4].metric("SNF 3-Day Rule Waiver", "Yes" if aco_data["SNF_Waiver"] == 1 else "No")
        cols_row2 = st.columns(3)
        cols_row2[0].metric("Assigned Beneficiaries", fmt_comma(aco_data["N_AB"]))
        cols_row2[1].metric("Quality Score", fmt_pct(aco_data["QualScore"]))
        cols_row2[2].metric("Earned Shared Savings/Loss", fmt_dollars(aco_data["EarnSaveLoss"]))
        # Facility Makeup
        st.markdown("### ACO Facility Makeup")
        st.write("Number of participating facilities by type (based on certified participant list and PECOS data).")
        facility_data = [
            {"Facility Type": "CAHs (Critical Access Hospitals)", "Count": fmt_comma(aco_data.get("N_CAH", "N/A"))},
            {"Facility Type": "FQHCs (Federally Qualified Health Centers)", "Count": fmt_comma(aco_data.get("N_FQHC", "N/A"))},
            {"Facility Type": "RHCs (Rural Health Clinics)", "Count": fmt_comma(aco_data.get("N_RHC", "N/A"))},
            {"Facility Type": "ETA Hospitals (Elected Teaching Amendment)", "Count": fmt_comma(aco_data.get("N_ETA", "N/A"))},
            {"Facility Type": "Short-term Acute Care Hospitals", "Count": fmt_comma(aco_data.get("N_Hosp", "N/A"))},
            {"Facility Type": "Other Facility Types", "Count": fmt_comma(aco_data.get("N_Fac_Other", "N/A"))},
        ]
        df_facility = pd.DataFrame(facility_data)
        st.dataframe(df_facility, use_container_width=True, hide_index=True)
        # Participating Provider Types
        st.markdown("### Participating Provider Types")
        st.write("Number of participating clinicians by type (reassigned billing rights to ACO participant).")
        provider_data = [
            {"Provider Type": "PCPs (Primary Care Physicians)", "Count": fmt_comma(aco_data.get("N_PCP", "N/A"))},
            {"Provider Type": "Specialists", "Count": fmt_comma(aco_data.get("N_Spec", "N/A"))},
            {"Provider Type": "Nurse Practitioners (NPs)", "Count": fmt_comma(aco_data.get("N_NP", "N/A"))},
            {"Provider Type": "Physician Assistants (PAs)", "Count": fmt_comma(aco_data.get("N_PA", "N/A"))},
            {"Provider Type": "Clinical Nurse Specialists (CNSs)", "Count": fmt_comma(aco_data.get("N_CNS", "N/A"))},
        ]
        df_provider = pd.DataFrame(provider_data)
        st.dataframe(df_provider, use_container_width=True, hide_index=True)
        # 2. Demographics Profile – smaller pies in rows
        st.markdown("### Demographics Profile")
        # Row 1: 3 pies
        pie_row1 = st.columns(3)
        with pie_row1[0]:
            enrollment_cols = {
                "ESRD": "N_AB_Year_ESRD_PY",
                "Disabled": "N_AB_Year_DIS_PY",
                "Aged Dual": "N_AB_Year_AGED_Dual_PY",
                "Aged Non-Dual": "N_AB_Year_AGED_NonDual_PY"
            }
            existing_enrollment = {k: v for k, v in enrollment_cols.items() if v in df.columns and pd.notna(aco_data[v])}
            if existing_enrollment:
                enrollment_df = pd.DataFrame({
                    "Enrollment Type": list(existing_enrollment.keys()),
                    "Person-Years": [aco_data[v] for v in existing_enrollment.values()]
                })
                fig_enrollment = px.pie(
                    enrollment_df,
                    values="Person-Years",
                    names="Enrollment Type",
                    title="Enrollment Status (Person-Years)",
                    hole=0.4
                )
                fig_enrollment.update_layout(template=PLOTLY_TEMPLATE, height=300, margin=dict(t=40, b=20, l=20, r=20))
                fig_enrollment.update_traces(textposition='inside', textinfo='percent+label')
                st.plotly_chart(fig_enrollment, use_container_width=True)
            else:
                st.info("No enrollment data.")
        with pie_row1[1]:
            align_cols = {
                "Voluntary Only": "N_Ben_VA_Only",
                "Claims-Based Only": "N_Ben_CBA_Only",
                "Both": "N_Ben_CBA_and_VA"
            }
            existing_align = {k: v for k, v in align_cols.items() if v in df.columns and pd.notna(aco_data[v])}
            if existing_align:
                align_df = pd.DataFrame({
                    "Alignment Type": list(existing_align.keys()),
                    "Beneficiaries": [aco_data[v] for v in existing_align.values()]
                })
                fig_align = px.pie(
                    align_df,
                    values="Beneficiaries",
                    names="Alignment Type",
                    title="Assignment Logic",
                    hole=0.4
                )
                fig_align.update_layout(template=PLOTLY_TEMPLATE, height=300, margin=dict(t=40, b=20, l=20, r=20))
                fig_align.update_traces(textposition='inside', textinfo='percent+label')
                st.plotly_chart(fig_align, use_container_width=True)
            else:
                st.info("No alignment data.")
        with pie_row1[2]:
            sex_cols = {
                "Male": "N_Ben_Male",
                "Female": "N_Ben_Female"
            }
            existing_sex = {k: v for k, v in sex_cols.items() if v in df.columns and pd.notna(aco_data[v])}
            if existing_sex:
                sex_df = pd.DataFrame({
                    "Sex": list(existing_sex.keys()),
                    "Beneficiaries": [aco_data[v] for v in existing_sex.values()]
                })
                fig_sex = px.pie(
                    sex_df,
                    values="Beneficiaries",
                    names="Sex",
                    title="Male:Female Ratio",
                    hole=0.4
                )
                fig_sex.update_layout(template=PLOTLY_TEMPLATE, height=300, margin=dict(t=40, b=20, l=20, r=20))
                fig_sex.update_traces(textposition='inside', textinfo='percent+label')
                st.plotly_chart(fig_sex, use_container_width=True)
            else:
                st.info("No sex data.")
        # Row 2: 2 pies (age + race)
        pie_row2 = st.columns(2)
        with pie_row2[0]:
            age_cols = {
                "Age 0-64": "N_Ben_Age_0_64",
                "Age 65-74": "N_Ben_Age_65_74",
                "Age 75-84": "N_Ben_Age_75_84",
                "Age 85+": "N_Ben_Age_85plus"
            }
            existing_age = {k: v for k, v in age_cols.items() if v in df.columns and pd.notna(aco_data[v])}
            if existing_age:
                age_df = pd.DataFrame({
                    "Age Group": list(existing_age.keys()),
                    "Beneficiaries": [aco_data[v] for v in existing_age.values()]
                })
                fig_age = px.pie(
                    age_df,
                    values="Beneficiaries",
                    names="Age Group",
                    title="Age Cohorts",
                    hole=0.4
                )
                fig_age.update_layout(template=PLOTLY_TEMPLATE, height=300, margin=dict(t=40, b=20, l=20, r=20))
                fig_age.update_traces(textposition='inside', textinfo='percent+label')
                st.plotly_chart(fig_age, use_container_width=True)
            else:
                st.info("No age data.")
        with pie_row2[1]:
            race_cols = {
                "White": "N_Ben_Race_White",
                "Black": "N_Ben_Race_Black",
                "Asian": "N_Ben_Race_Asian",
                "Hispanic": "N_Ben_Race_Hisp",
                "Native American": "N_Ben_Race_Native",
                "Other": "N_Ben_Race_Other",
                "Unknown": "N_Ben_Race_Unknown"
            }
            existing_race = {k: v for k, v in race_cols.items() if v in df.columns and pd.notna(aco_data[v])}
            if existing_race:
                race_df = pd.DataFrame({
                    "Race": list(existing_race.keys()),
                    "Beneficiaries": [aco_data[v] for v in existing_race.values()]
                })
                fig_race = px.pie(
                    race_df,
                    values="Beneficiaries",
                    names="Race",
                    title="Race Breakdown",
                    hole=0.4
                )
                fig_race.update_layout(template=PLOTLY_TEMPLATE, height=300, margin=dict(t=40, b=20, l=20, r=20))
                fig_race.update_traces(textposition='inside', textinfo='percent+label')
                st.plotly_chart(fig_race, use_container_width=True)
            else:
                st.info("No race data.")
        # Risk Adjustment section
        st.markdown("### Risk Adjustment")
        risk_cols = {
            "ESRD": ["CMS_HCC_RiskScore_ESRD_BY1", "CMS_HCC_RiskScore_ESRD_BY2", "CMS_HCC_RiskScore_ESRD_BY3", "CMS_HCC_RiskScore_ESRD_PY"],
            "Disabled": ["CMS_HCC_RiskScore_DIS_BY1", "CMS_HCC_RiskScore_DIS_BY2", "CMS_HCC_RiskScore_DIS_BY3", "CMS_HCC_RiskScore_DIS_PY"],
            "Aged Dual": ["CMS_HCC_RiskScore_AGDU_BY1", "CMS_HCC_RiskScore_AGDU_BY2", "CMS_HCC_RiskScore_AGDU_BY3", "CMS_HCC_RiskScore_AGDU_PY"],
            "Aged Non-Dual": ["CMS_HCC_RiskScore_AGND_BY1", "CMS_HCC_RiskScore_AGND_BY2", "CMS_HCC_RiskScore_AGND_BY3", "CMS_HCC_RiskScore_AGND_PY"]
        }
        risk_df = pd.DataFrame({
            "Category": list(risk_cols.keys())
        })
        periods = ["BY1", "BY2", "BY3", "PY"]
        for i, period in enumerate(periods):
            risk_df[period] = [aco_data[risk_cols[cat][i]] if risk_cols[cat][i] in df.columns else np.nan for cat in risk_cols]
        risk_df["Track Avg (PY)"] = [track_avg.get(risk_cols[cat][3], np.nan) for cat in risk_cols]
        st.dataframe(risk_df, use_container_width=True, hide_index=True)
        # Use pre-calculated weighted risk scores from df
        weighted_risk_py = aco_data.get("weighted_risk_py", np.nan)
        if pd.notna(weighted_risk_py):
            st.metric("Overall Weighted Risk Score (PY)", f"{weighted_risk_py:.3f}")
        else:
            st.metric("Overall Weighted Risk Score (PY)", "N/A")
        # 5. Financial & Performance Roll-Up (preserved)
        st.markdown("### Financial & Performance Roll-Up")
        with st.expander("Expand to see detailed calculation path", expanded=False):
            # Summary level table (visible inside expander)
            summary_data = {
                "Metric": [
                    "Total Benchmark",
                    "Total Expenditures",
                    "Raw Savings/Loss",
                    "Savings Rate",
                    "Earned Shared Savings/Loss"
                ],
                "Aggregate ($)": [
                    fmt_dollars(aco_data["ABtotBnchmk"]),
                    fmt_dollars(aco_data["ABtotExp"]),
                    fmt_dollars(aco_data["BnchmkMinExp"]),
                    fmt_pct(aco_data["Sav_rate"]),
                    fmt_dollars(aco_data["EarnSaveLoss"])
                ],
                "Per Beneficiary (Annual)": [
                    fmt_dollars(aco_data["ABtotBnchmk"] / aco_data["N_AB_Year_PY"] if aco_data["N_AB_Year_PY"] > 0 else 0),
                    fmt_dollars(aco_data["ABtotExp"] / aco_data["N_AB_Year_PY"] if aco_data["N_AB_Year_PY"] > 0 else 0),
                    fmt_dollars(aco_data["BnchmkMinExp"] / aco_data["N_AB_Year_PY"] if aco_data["N_AB_Year_PY"] > 0 else 0),
                    "-",
                    fmt_dollars(aco_data["EarnSaveLoss"] / aco_data["N_AB_Year_PY"] if aco_data["N_AB_Year_PY"] > 0 else 0)
                ],
                "PMPM": [
                    fmt_dollars(aco_data["ABtotBnchmk"] / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    fmt_dollars(aco_data["ABtotExp"] / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    fmt_dollars(aco_data["BnchmkMinExp"] / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    "-",
                    fmt_dollars(aco_data["EarnSaveLoss"] / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2)
                ],
                "vs Track Avg PMPM": [
                    fmt_dollars((aco_data["ABtotBnchmk"] - track_avg["ABtotBnchmk"]) / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    fmt_dollars((aco_data["ABtotExp"] - track_avg["ABtotExp"]) / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    fmt_dollars((aco_data["BnchmkMinExp"] - track_avg["BnchmkMinExp"]) / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2),
                    "-",
                    fmt_dollars((aco_data["EarnSaveLoss"] - track_avg["EarnSaveLoss"]) / aco_data["N_AB_Year_PY"] / 12 if aco_data["N_AB_Year_PY"] > 0 else 0, 2)
                ]
            }
            summary_df = pd.DataFrame(summary_data)
            st.dataframe(summary_df, use_container_width=True, hide_index=True)
            # Export button
            csv = aco_data.to_csv(index=True).encode('utf-8')
            st.download_button(
                label="Export Full ACO Data (CSV)",
                data=csv,
                file_name=f"{aco_data['ACO_Name']}_2023.csv",
                mime="text/csv"
            )

            # ── PPTX Export Button ───────────────────────────────────────────────
            if st.button("Export Full Report (PPTX)"):
                pptx_bytes = generate_pptx_report(aco_data, df, track_avg)
                st.download_button(
                    label="Download PPTX Now",
                    data=pptx_bytes,
                    file_name=f"MSSP_PY2023_Report_{aco_data.get('ACO_ID', 'SelectedACO')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"pptx_{aco_data.get('ACO_ID')}"
                )

            # Nested expanders for details
            with st.expander("Total Benchmark Components"):
                st.write("**Calculation**: Historical Benchmark + Regional/Trend Adjustment = Final Benchmark")
                st.write(f"- Historical Benchmark: {fmt_dollars(aco_data.get('HistBnchmk', np.nan))}")
                adjustment = aco_data.get('UpdatedBnchmk', np.nan) - aco_data.get('HistBnchmk', np.nan) if pd.notna(aco_data.get('HistBnchmk')) and pd.notna(aco_data.get('UpdatedBnchmk')) else np.nan
                st.write(f"- Regional/Trend Adjustment: {fmt_dollars(adjustment)}")
                st.write(f"- Final Benchmark Used: {fmt_dollars(aco_data['ABtotBnchmk'])}")
                st.info("Note: Full historical baseline years and regional factors are in CMS historical benchmark files (not in this PUF).")
            with st.expander("Total Expenditures Breakdown"):
                st.write("**Breakdown by major service category** (annual aggregate)")
                capann_cols = {
                    "Inpatient (Total)": "CapAnn_INP_All",
                    "Outpatient Dept": "CapAnn_OPD",
                    "Physician/Supplier": "CapAnn_PB",
                    "SNF": "CapAnn_SNF",
                    "Home Health": "CapAnn_HHA",
                    "Hospice": "CapAnn_HSP",
                    "DME": "CapAnn_DME"
                }
                existing_cols = {k: v for k, v in capann_cols.items() if v in df.columns}
                if existing_cols:
                    # Calculate component aggregates using person-years
                    exp_df = pd.DataFrame({
                        "Category": list(existing_cols.keys()),
                        "Aggregate ($)": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] if pd.notna(aco_data.get(v)) else np.nan) for v in existing_cols.values()],
                        "Per Beneficiary": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] / aco_data["N_AB"] if pd.notna(aco_data.get(v)) and aco_data["N_AB"] > 0 else 0) for v in existing_cols.values()],
                        "PMPM": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] / aco_data["N_AB"] / 12 if pd.notna(aco_data.get(v)) and aco_data["N_AB"] > 0 else 0, 2) for v in existing_cols.values()]
                    })
                    # Sum of components (aggregate)
                    component_sum_agg = sum(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] for v in existing_cols.values() if pd.notna(aco_data.get(v)))
                    # Truncation adjustment (aggregate)
                    total_exp = aco_data["ABtotExp"] if pd.notna(aco_data["ABtotExp"]) else 0
                    truncation_adj_agg = total_exp - component_sum_agg
                    # Subtotal and adjustment rows
                    exp_df.loc[len(exp_df)] = ["Subtotal (Components)", fmt_dollars(component_sum_agg), fmt_dollars(component_sum_agg / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(component_sum_agg / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                    exp_df.loc[len(exp_df)] = ["Truncation Adjustment", fmt_dollars(truncation_adj_agg), fmt_dollars(truncation_adj_agg / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(truncation_adj_agg / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                    exp_df.loc[len(exp_df)] = ["Total Expenditures", fmt_dollars(total_exp), fmt_dollars(total_exp / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(total_exp / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                    st.dataframe(exp_df, use_container_width=True, hide_index=True)
                    # Note on truncation
                    st.info("**Note**: Small discrepancies between total and component sum are due to CMS truncation/rounding of per capita values to whole dollars. Adjustment row reconciles to official total (ABtotExp). Per Beneficiary uses total assigned beneficiaries (N_AB).")
                    # Nested expander for Inpatient (Total)
                    with st.expander("Inpatient (Total) Breakdown"):
                        inpatient_cols = {
                            "Short-Term": "CapAnn_INP_S_trm",
                            "Long-Term": "CapAnn_INP_L_trm",
                            "Rehab": "CapAnn_INP_Rehab",
                            "Psych": "CapAnn_INP_Psych"
                        }
                        existing_inpatient = {k: v for k, v in inpatient_cols.items() if v in df.columns}
                        if existing_inpatient:
                            inp_df = pd.DataFrame({
                                "Sub-Category": list(existing_inpatient.keys()),
                                "Aggregate ($)": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] if pd.notna(aco_data.get(v)) else np.nan) for v in existing_inpatient.values()],
                                "Per Beneficiary": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] / aco_data["N_AB"] if pd.notna(aco_data.get(v)) and aco_data["N_AB"] > 0 else 0) for v in existing_inpatient.values()],
                                "PMPM": [fmt_dollars(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] / aco_data["N_AB"] / 12 if pd.notna(aco_data.get(v)) and aco_data["N_AB"] > 0 else 0, 2) for v in existing_inpatient.values()]
                            })
                            # Sum of inpatient components (aggregate)
                            component_sum_agg = sum(aco_data.get(v, 0) * aco_data["N_AB_Year_PY"] for v in existing_inpatient.values() if pd.notna(aco_data.get(v)))
                            # Truncation adjustment (aggregate)
                            inp_total = aco_data["CapAnn_INP_All"] * aco_data["N_AB_Year_PY"] if pd.notna(aco_data["CapAnn_INP_All"]) else 0
                            truncation_adj_agg = inp_total - component_sum_agg
                            # Subtotal and adjustment rows
                            inp_df.loc[len(inp_df)] = ["Subtotal (Components)", fmt_dollars(component_sum_agg), fmt_dollars(component_sum_agg / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(component_sum_agg / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                            inp_df.loc[len(inp_df)] = ["Truncation Adjustment", fmt_dollars(truncation_adj_agg), fmt_dollars(truncation_adj_agg / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(truncation_adj_agg / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                            inp_df.loc[len(inp_df)] = ["Total Inpatient", fmt_dollars(inp_total), fmt_dollars(inp_total / aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0), fmt_dollars(inp_total / aco_data["N_AB"] / 12 if aco_data["N_AB"] > 0 else 0, 2)]
                            st.dataframe(inp_df, use_container_width=True, hide_index=True)
                        else:
                            st.info("No detailed inpatient breakdown columns available.")
                else:
                    st.info("No detailed expenditure breakdown columns available.")
            with st.expander("Raw Savings/Loss Details"):
                st.write("**Calculation**: Total Benchmark - Total Expenditures")
                st.write(f"- Benchmark: {fmt_dollars(aco_data['ABtotBnchmk'])}")
                st.write(f"- Expenditures: {fmt_dollars(aco_data['ABtotExp'])}")
                st.write(f"- Raw Savings/Loss: {fmt_dollars(aco_data['BnchmkMinExp'])}")
                st.write(f"- Minimum Savings Rate (MSR): {fmt_pct(aco_data['MinSavPerc'])} (threshold for sharing eligibility)")
            with st.expander("Savings Rate Details"):
                st.write("**Calculation**: Raw Savings/Loss / Total Benchmark")
                st.write(f"- Raw Savings/Loss: {fmt_dollars(aco_data['BnchmkMinExp'])}")
                st.write(f"- Total Benchmark: {fmt_dollars(aco_data['ABtotBnchmk'])}")
                st.write(f"- Savings Rate: {fmt_pct(aco_data['Sav_rate'])}")
            with st.expander("Earned Shared Savings/Loss Details"):
                st.write("**Calculation**: Generated Savings/Loss × Quality Sharing Rate - Adjustments")
                st.write(f"- Generated Savings/Loss: {fmt_dollars(aco_data['GenSaveLoss'])}")
                st.write(f"- Final Quality Sharing Rate: {fmt_pct(aco_data['FinalShareRate'])} (based on quality performance)")
                st.write(f"- EUC Adjustment (if applicable): {fmt_dollars(aco_data['DisAdj'])}")
                st.write(f"- Mid-Year Termination Proration: {'Yes (prorated)' if aco_data['Impact_Mid_Year_Termination'] == 1 else 'No'}")
                st.write(f"- Sequestration / Payment Limit: 2% federal reduction applied per policy (not ACO-specific in PUF)")
                st.write(f"- Final Earned Shared Savings/Loss: {fmt_dollars(aco_data['EarnSaveLoss'])}")
        # 5. Utilization Breakdown Chart with Physician/Supplier added + % of Total Spend pie
        st.markdown("### Utilization Breakdown")
        capann_cols = {
            "Inpatient (Total)": "CapAnn_INP_All",
            "Outpatient Dept": "CapAnn_OPD",
            "Physician/Supplier": "CapAnn_PB",
            "SNF": "CapAnn_SNF",
            "Home Health": "CapAnn_HHA",
            "Hospice": "CapAnn_HSP",
            "DME": "CapAnn_DME"
        }
        existing_cols = {k: v for k, v in capann_cols.items() if v in df.columns}
        if existing_cols:
            # Existing per capita bar chart
            util_melt = pd.DataFrame({
                "Category": list(existing_cols.keys()),
                "This ACO": [aco_data.get(v, 0) for v in existing_cols.values()],
                "Track Average": [track_avg.get(v, 0) for v in existing_cols.values()]
            }).melt(id_vars="Category", var_name="Type", value_name="Per Capita $")
            fig_util = px.bar(
                util_melt,
                x="Category",
                y="Per Capita $",
                color="Type",
                barmode="group",
                title="Per Capita Expenditures by Service Category – This ACO vs Track Average",
                height=500,
                category_orders={"Category": list(capann_cols.keys())}
            )
            fig_util.update_layout(template=PLOTLY_TEMPLATE, xaxis_tickangle=-45)
            # New: % of Total Spend pie (using This ACO values)
            total_util = sum(aco_data.get(v, 0) for v in existing_cols.values() if pd.notna(aco_data.get(v)))
            if total_util > 0:
                percent_df = pd.DataFrame({
                    "Category": list(existing_cols.keys()),
                    "% of Total": [(aco_data.get(v, 0) / total_util * 100) for v in existing_cols.values()]
                })
                fig_percent = px.pie(
                    percent_df,
                    values="% of Total",
                    names="Category",
                    title="% of Total Per Capita Utilization – This ACO",
                    hole=0.4
                )
                fig_percent.update_traces(textposition='inside', textinfo='percent+label')
                fig_percent.update_layout(template=PLOTLY_TEMPLATE, height=400)
                # Side-by-side layout
                col_left, col_right = st.columns(2)
                with col_left:
                    st.plotly_chart(fig_util, use_container_width=True)
                with col_right:
                    st.plotly_chart(fig_percent, use_container_width=True)
            else:
                st.plotly_chart(fig_util, use_container_width=True)
                st.info("Total utilization is zero or NaN – % pie not shown.")
        else:
            st.info("No utilization expenditure columns available for this ACO.")
        # Inpatient Derived Metrics Table
        st.markdown("### Inpatient Utilization Metrics")
        st.write("Derived metrics for inpatient discharges by type (per 1,000 person-years unless noted).")
        inpatient_types = [
            {"Name": "Inpatient Total", "adm_col": "ADM", "capann_col": "CapAnn_INP_All"},
            {"Name": "Inpatient - Short Term Acute Care", "adm_col": "ADM_S_Trm", "capann_col": "CapAnn_INP_S_trm"},
            {"Name": "Inpatient - Long Term Care", "adm_col": "ADM_L_Trm", "capann_col": "CapAnn_INP_L_trm"},
            {"Name": "Inpatient Rehab Facility", "adm_col": "ADM_Rehab", "capann_col": "CapAnn_INP_Rehab"},
            {"Name": "Inpatient Psychiatric Facility", "adm_col": "ADM_Psych", "capann_col": "CapAnn_INP_Psych"},
        ]
        inpatient_table = []
        for typ in inpatient_types:
            adm_col = typ["adm_col"]
            capann_col = typ["capann_col"]
            if adm_col not in df.columns or capann_col not in df.columns:
                continue
            admits_per_1000 = aco_data.get(adm_col, 0)
            aggregate_admits = admits_per_1000 * (aco_data.get("N_AB_Year_PY", 0) / 1000) if aco_data.get("N_AB_Year_PY", 0) > 0 else 0
            total_exp = aco_data.get(capann_col, 0) * aco_data.get("N_AB_Year_PY", 0)
            cost_per_admit = total_exp / aggregate_admits if aggregate_admits > 0 else 0
            cost_per_benef = total_exp / aco_data.get("N_AB_Year_PY", 1) if aco_data.get("N_AB_Year_PY", 1) > 0 else 0
            pmpm = cost_per_benef / 12
            inpatient_table.append({
                "Inpatient Type": typ["Name"],
                "Admits Per 1,000 Beneficiaries": f"{admits_per_1000:.0f}",
                "Aggregate Admits": f"{aggregate_admits:,.0f}",
                "Total Expenditures": fmt_dollars(total_exp, 0),
                "Cost/Admit": fmt_dollars(cost_per_admit, 2),
                "Cost/Beneficiary": fmt_dollars(cost_per_benef, 0),
                "PMPM": fmt_dollars(pmpm, 2)
            })
        if inpatient_table:
            df_inpatient = pd.DataFrame(inpatient_table)
            st.dataframe(df_inpatient, use_container_width=True, hide_index=True)
        else:
            st.info("No inpatient discharge columns found in data. Check CSV for ADM, ADM_S_Trm, CapAnn_INP_All, etc.")
        # SNF Utilization Metrics
        st.markdown("### SNF Utilization Metrics")
        st.write("Skilled Nursing Facility metrics vs average across all ACOs.")
        snf_table = []
        # 1. SNF Admissions per 1,000 Person-Years
        if "P_SNF_ADM" in df.columns:
            this_aco_adm = aco_data.get("P_SNF_ADM", np.nan)
            all_aco_avg_adm = df["P_SNF_ADM"].mean(skipna=True)
            pct_diff_adm = (this_aco_adm - all_aco_avg_adm) / all_aco_avg_adm if pd.notna(all_aco_avg_adm) and all_aco_avg_adm != 0 else np.nan
            deviation_adm = f"{pct_diff_adm*100:.1f}%" if pd.notna(pct_diff_adm) else "N/A"
            aggregate_adm = this_aco_adm * (aco_data.get("N_AB_Year_PY", 0) / 1000) if aco_data.get("N_AB_Year_PY", 0) > 0 else 0
            snf_table.append({
                "Metric": "SNF Admissions per 1,000 Person-Years",
                "This ACO": f"{this_aco_adm:.1f}" if pd.notna(this_aco_adm) else "N/A",
                "All ACOs Avg": f"{all_aco_avg_adm:.1f}" if pd.notna(all_aco_avg_adm) else "N/A",
                "% Difference": deviation_adm
            })
        # 2. Aggregate SNF Admissions
        if "P_SNF_ADM" in df.columns:
            aggregate_adm_this = this_aco_adm * (aco_data.get("N_AB_Year_PY", 0) / 1000) if aco_data.get("N_AB_Year_PY", 0) > 0 else 0
            aggregate_adm_avg = all_aco_avg_adm * (df["N_AB_Year_PY"].mean(skipna=True) / 1000) if pd.notna(all_aco_avg_adm) else np.nan
            pct_diff_agg_adm = (aggregate_adm_this - aggregate_adm_avg) / aggregate_adm_avg if pd.notna(aggregate_adm_avg) and aggregate_adm_avg != 0 else np.nan
            deviation_agg_adm = f"{pct_diff_agg_adm*100:.1f}%" if pd.notna(pct_diff_agg_adm) else "N/A"
            snf_table.append({
                "Metric": "Aggregate SNF Admissions",
                "This ACO": f"{aggregate_adm_this:,.0f}" if pd.notna(aggregate_adm_this) else "N/A",
                "All ACOs Avg": f"{aggregate_adm_avg:,.0f}" if pd.notna(aggregate_adm_avg) else "N/A",
                "% Difference": deviation_agg_adm
            })
        # 3. Total SNF Expenditures
        if "CapAnn_SNF" in df.columns:
            total_exp_this = aco_data.get("CapAnn_SNF", 0) * aco_data.get("N_AB_Year_PY", 0)
            total_exp_avg = df["CapAnn_SNF"].mean(skipna=True) * df["N_AB_Year_PY"].mean(skipna=True)
            pct_diff_exp = (total_exp_this - total_exp_avg) / total_exp_avg if pd.notna(total_exp_avg) and total_exp_avg != 0 else np.nan
            deviation_exp = f"{pct_diff_exp*100:.1f}%" if pd.notna(pct_diff_exp) else "N/A"
            snf_table.append({
                "Metric": "Total SNF Expenditures",
                "This ACO": fmt_dollars(total_exp_this, 0) if pd.notna(total_exp_this) else "N/A",
                "All ACOs Avg": fmt_dollars(total_exp_avg, 0) if pd.notna(total_exp_avg) else "N/A",
                "% Difference": deviation_exp
            })
        # 4. Average SNF LOS (days)
        if "SNF_LOS" in df.columns:
            los_this = aco_data.get("SNF_LOS", np.nan)
            los_avg = df["SNF_LOS"].mean(skipna=True)
            pct_diff_los = (los_this - los_avg) / los_avg if pd.notna(los_avg) and los_avg != 0 else np.nan
            deviation_los = f"{pct_diff_los*100:.1f}%" if pd.notna(pct_diff_los) else "N/A"
            snf_table.append({
                "Metric": "Average SNF LOS (days)",
                "This ACO": f"{los_this:.1f}" if pd.notna(los_this) else "N/A",
                "All ACOs Avg": f"{los_avg:.1f}" if pd.notna(los_avg) else "N/A",
                "% Difference": deviation_los
            })
        # 5. Cost per SNF Day
        if "CapAnn_SNF" in df.columns and "SNF_LOS" in df.columns:
            snf_days_this = (aco_data.get("SNF_LOS", 0) * aggregate_adm_this) if 'aggregate_adm_this' in locals() else 0
            cost_per_day_this = total_exp_this / snf_days_this if snf_days_this > 0 else 0
            snf_days_avg = (los_avg * aggregate_adm_avg) if 'los_avg' in locals() and 'aggregate_adm_avg' in locals() else 0
            cost_per_day_avg = total_exp_avg / snf_days_avg if snf_days_avg > 0 else 0
            pct_diff_day = (cost_per_day_this - cost_per_day_avg) / cost_per_day_avg if pd.notna(cost_per_day_avg) and cost_per_day_avg != 0 else np.nan
            deviation_day = f"{pct_diff_day*100:.1f}%" if pd.notna(pct_diff_day) else "N/A"
            snf_table.append({
                "Metric": "Cost per SNF Day",
                "This ACO": fmt_dollars(cost_per_day_this, 0) if pd.notna(cost_per_day_this) else "N/A",
                "All ACOs Avg": fmt_dollars(cost_per_day_avg, 0) if pd.notna(cost_per_day_avg) else "N/A",
                "% Difference": deviation_day
            })
        # 6. Cost per SNF Admission
        if 'total_exp_this' in locals() and 'aggregate_adm_this' in locals():
            cost_per_adm_this = total_exp_this / aggregate_adm_this if aggregate_adm_this > 0 else 0
            cost_per_adm_avg = total_exp_avg / aggregate_adm_avg if aggregate_adm_avg > 0 else 0
            pct_diff_adm_cost = (cost_per_adm_this - cost_per_adm_avg) / cost_per_adm_avg if pd.notna(cost_per_adm_avg) and cost_per_adm_avg != 0 else np.nan
            deviation_adm_cost = f"{pct_diff_adm_cost*100:.1f}%" if pd.notna(pct_diff_adm_cost) else "N/A"
            snf_table.append({
                "Metric": "Cost per SNF Admission",
                "This ACO": fmt_dollars(cost_per_adm_this, 2) if pd.notna(cost_per_adm_this) else "N/A",
                "All ACOs Avg": fmt_dollars(cost_per_adm_avg, 2) if pd.notna(cost_per_adm_avg) else "N/A",
                "% Difference": deviation_adm_cost
            })
        # 7. SNF Cost per Beneficiary & PMPM
        if 'total_exp_this' in locals():
            cost_per_benef_this = total_exp_this / aco_data.get("N_AB", 1) if aco_data.get("N_AB", 1) > 0 else 0
            pmpm_this = cost_per_benef_this / 12
            cost_per_benef_avg = total_exp_avg / df["N_AB"].mean(skipna=True) if df["N_AB"].mean(skipna=True) > 0 else 0
            pmpm_avg = cost_per_benef_avg / 12
            pct_diff_benef = (cost_per_benef_this - cost_per_benef_avg) / cost_per_benef_avg if pd.notna(cost_per_benef_avg) and cost_per_benef_avg != 0 else np.nan
            deviation_benef = f"{pct_diff_benef*100:.1f}%" if pd.notna(pct_diff_benef) else "N/A"
            snf_table.append({
                "Metric": "SNF Cost per Beneficiary",
                "This ACO": fmt_dollars(cost_per_benef_this, 0) if pd.notna(cost_per_benef_this) else "N/A",
                "All ACOs Avg": fmt_dollars(cost_per_benef_avg, 0) if pd.notna(cost_per_benef_avg) else "N/A",
                "% Difference": deviation_benef
            })
            snf_table.append({
                "Metric": "SNF PMPM",
                "This ACO": fmt_dollars(pmpm_this, 2) if pd.notna(pmpm_this) else "N/A",
                "All ACOs Avg": fmt_dollars(pmpm_avg, 2) if pd.notna(pmpm_avg) else "N/A",
                "% Difference": f"{(pmpm_this - pmpm_avg) / pmpm_avg * 100:.1f}%" if pd.notna(pmpm_avg) and pmpm_avg != 0 else "N/A"
            })
        if snf_table:
            df_snf = pd.DataFrame(snf_table)
            st.dataframe(df_snf, use_container_width=True, hide_index=True)
        else:
            st.info("No SNF columns found (P_SNF_ADM, CapAnn_SNF, SNF_LOS missing).")
        # Emergency Department (ED) Utilization Comparison
        st.markdown("### Emergency Department (ED) Utilization")
        st.write("ED visits per 1,000 person-years and % leading to hospitalization vs average across all ACOs. Deviations >20% highlighted.")
        ed_data = []
        # ED Visits per 1,000 person-years
        if "P_EDV_Vis" in df.columns:
            this_aco_ed_vis = aco_data.get("P_EDV_Vis", np.nan)
            all_aco_avg_ed_vis = df["P_EDV_Vis"].mean() if "P_EDV_Vis" in df.columns else np.nan
            pct_diff_vis = (this_aco_ed_vis - all_aco_avg_ed_vis) / all_aco_avg_ed_vis if pd.notna(all_aco_avg_ed_vis) and all_aco_avg_ed_vis != 0 else np.nan
            deviation_vis = f"{pct_diff_vis*100:.1f}%" if pd.notna(pct_diff_vis) else "N/A"
            alert_vis = "High" if pct_diff_vis > 0.20 else "Low" if pct_diff_vis < -0.20 else ""
            ed_data.append({
                "Metric": "ED Visits per 1,000 Person-Years",
                "This ACO": f"{this_aco_ed_vis:.1f}" if pd.notna(this_aco_ed_vis) else "N/A",
                "All ACOs Avg": f"{all_aco_avg_ed_vis:.1f}" if pd.notna(all_aco_avg_ed_vis) else "N/A",
                "% Difference": deviation_vis,
                "Alert": alert_vis
            })
        # % of ED Visits Leading to Hospitalization
        if "P_EDV_Vis_HOSP" in df.columns:
            this_aco_ed_vis = aco_data.get("P_EDV_Vis_HOSP", np.nan)
            all_aco_avg_ed_vis = df["P_EDV_Vis_HOSP"].mean() if "P_EDV_Vis_HOSP" in df.columns else np.nan
            pct_diff_vis = (this_aco_ed_vis - all_aco_avg_ed_vis) / all_aco_avg_ed_vis if pd.notna(all_aco_avg_ed_vis) and all_aco_avg_ed_vis != 0 else np.nan
            deviation_vis = f"{pct_diff_vis*100:.1f}%" if pd.notna(pct_diff_vis) else "N/A"
            alert_vis = "High" if pct_diff_vis > 0.20 else "Low" if pct_diff_vis < -0.20 else ""
            ed_data.append({
                "Metric": "ED Visits Leading to Hospitalization per 1,000 Person-Years",
                "This ACO": f"{this_aco_ed_vis:.1f}" if pd.notna(this_aco_ed_vis) else "N/A",
                "All ACOs Avg": f"{all_aco_avg_ed_vis:.1f}" if pd.notna(all_aco_avg_ed_vis) else "N/A",
                "% Difference": deviation_vis,
                "Alert": alert_vis
            })
        if ed_data:
            df_ed = pd.DataFrame(ed_data)
            # Highlight alerts
            def highlight_alert(row):
                if row["Alert"] == "High":
                    return ['background-color: #ffcccc'] * len(row)
                elif row["Alert"] == "Low":
                    return ['background-color: #ccffcc'] * len(row)
                return [''] * len(row)
            st.dataframe(
                df_ed.style.apply(highlight_alert, axis=1),
                use_container_width=True,
                hide_index=True
            )
        else:
            st.info("ED columns (P_EDV_Vis and/or P_EDV_Vis_Hosp) not found in data.")
        # Primary Care Services View
        st.markdown("### Primary Care Services")
        st.write("E&M visits per 1,000 person-years by provider type vs average across all ACOs.")
        pcp_metrics = [
            {"Name": "Total E&M Visits per 1,000 Person-Years", "col": "P_EM_Total"},
            {"Name": "PCP E&M Visits per 1,000 Person-Years", "col": "P_EM_PCP_Vis"},
            {"Name": "Specialist E&M Visits per 1,000 Person-Years", "col": "P_EM_SP_Vis"},
            {"Name": "Nurse Visits per 1,000 Person-Years", "col": "P_Nurse_Vis"},
            {"Name": "FQHC/RHC Visits per 1,000 Person-Years", "col": "P_FQHC_RHC_Vis"},
        ]
        pcp_table = []
        for m in pcp_metrics:
            col = m["col"]
            if col not in df.columns:
                continue
            # Force numeric (handles suppression like *, -, empty)
            df_numeric = pd.to_numeric(df[col], errors='coerce')
            this_aco_val = pd.to_numeric(aco_data.get(col, np.nan), errors='coerce')
            all_aco_avg = df_numeric.mean(skipna=True)
            pct_diff = (this_aco_val - all_aco_avg) / all_aco_avg if pd.notna(all_aco_avg) and all_aco_avg != 0 else np.nan
            deviation_str = f"{pct_diff*100:.1f}%" if pd.notna(pct_diff) else "N/A"
            # Format with commas for thousands
            this_aco_fmt = f"{this_aco_val:,.1f}" if pd.notna(this_aco_val) else "N/A"
            all_aco_fmt = f"{all_aco_avg:,.1f}" if pd.notna(all_aco_avg) else "N/A"
            pcp_table.append({
                "Metric": m["Name"],
                "This ACO": this_aco_fmt,
                "All ACOs Avg": all_aco_fmt,
                "% Difference": deviation_str
            })
        if pcp_table:
            df_pcp = pd.DataFrame(pcp_table)
            st.dataframe(df_pcp, use_container_width=True, hide_index=True)
        else:
            st.info("No primary care columns found (P_EM_Total, P_EM_PCP_Vis, etc. missing).")
        # PCP Services vs Raw Savings Scatter Plot with Callout (ACO Name)
        st.markdown("### PCP Services vs Raw Savings")
        st.write("Scatter plot across all ACOs: PCP E&M visits per 1,000 person-years vs raw savings. Bubble size = Assigned Beneficiaries; color = Quality Score (green = high, red = low). Selected ACO highlighted with callout.")
        if all(col in df.columns for col in ["P_EM_PCP_Vis", "BnchmkMinExp", "QualScore", "N_AB", "ACO_Name", "ACO_ID"]):
            plot_df = df.copy()
            plot_df = plot_df[plot_df["P_EM_PCP_Vis"].notna() & plot_df["BnchmkMinExp"].notna() & plot_df["QualScore"].notna() & plot_df["N_AB"].notna()]
            fig_scat = px.scatter(
                plot_df,
                x="P_EM_PCP_Vis",
                y="BnchmkMinExp",
                size="N_AB",
                color="QualScore",
                color_continuous_scale="RdYlGn", # Red (low) → Yellow → Green (high)
                trendline="ols",
                trendline_scope="overall",
                trendline_color_override="#A5D8FF",
                hover_name="ACO_Name",
                hover_data={
                    "P_EM_PCP_Vis": ":.1f",
                    "BnchmkMinExp": ":$,.0f",
                    "QualScore": ":.1f",
                    "N_AB": ":,",
                    "Current_Track": True
                },
                title="PCP E&M Visits per 1,000 vs Raw Savings<br>Bubble size = Assigned Beneficiaries; Color = Quality Score",
                labels={
                    "P_EM_PCP_Vis": "PCP E&M Visits per 1,000 Person-Years",
                    "BnchmkMinExp": "Raw Savings ($)"
                }
            )
            # Add callout annotation with actual ACO name
            selected_row = plot_df[plot_df["ACO_ID"] == aco_data["ACO_ID"]]
            if not selected_row.empty:
                x_val = selected_row["P_EM_PCP_Vis"].iloc[0]
                y_val = selected_row["BnchmkMinExp"].iloc[0]
                aco_name = selected_row["ACO_Name"].iloc[0]
                fig_scat.add_annotation(
                    x=x_val,
                    y=y_val,
                    text=f"<b>{aco_name}</b>",
                    showarrow=True,
                    arrowhead=2,
                    ax=60, # Arrow x-offset — adjust if needed
                    ay=-60, # Arrow y-offset — adjust if needed
                    font=dict(size=14, color="#000000"),
                    bgcolor="#FFFFFF",
                    bordercolor="#000000",
                    borderwidth=2,
                    arrowwidth=2,
                    arrowcolor="#000000"
                )
            fig_scat.update_traces(marker=dict(opacity=0.7, line=dict(width=0.5)))
            fig_scat.update_layout(template=PLOTLY_TEMPLATE, height=600)
            st.plotly_chart(fig_scat, use_container_width=True)
        else:
            st.info("Missing columns for scatter plot (P_EM_PCP_Vis, BnchmkMinExp, QualScore, N_AB, ACO_ID).")
        # PCP Services vs Raw Savings Per Beneficiary Scatter Plot (Fallback Size)
        st.markdown("### PCP Services vs Raw Savings Per Beneficiary")
        st.write("Scatter plot across all ACOs: PCP E&M visits per 1,000 person-years vs raw savings per beneficiary. Selected ACO bubble larger (scaled by risk score). Trend line added.")
        if all(col in df.columns for col in ["P_EM_PCP_Vis", "BnchmkMinExp", "N_AB_Year_PY", "ACO_Name", "ACO_ID"]):
            plot_df = df.copy()
            plot_df = plot_df[plot_df["P_EM_PCP_Vis"].notna() & plot_df["BnchmkMinExp"].notna() & plot_df["N_AB_Year_PY"].notna() & (plot_df["N_AB_Year_PY"] > 0)]
            # Calculate raw savings per beneficiary
            plot_df["Raw_Savings_Per_Benef"] = plot_df["BnchmkMinExp"] / plot_df["N_AB_Year_PY"]
            # Create list for bubble size: every ACO gets its own risk score scaled
            size = [np.nan_to_num(row["weighted_risk_py"], nan=0.1) * 100 for _, row in plot_df.iterrows()]
            fig_scat_per = px.scatter(
                plot_df,
                x="P_EM_PCP_Vis",
                y="Raw_Savings_Per_Benef",
                size=size, # Now uses the safe list
                trendline="ols",
                trendline_scope="overall",
                trendline_color_override="#A5D8FF",
                hover_name="ACO_Name",
                hover_data={
                    "P_EM_PCP_Vis": ":.1f",
                    "Raw_Savings_Per_Benef": ":$,.2f",
                    "N_AB_Year_PY": ":,",
                    "weighted_risk_py": ":.3f",
                    "Current_Track": True
                },
                title="PCP E&M Visits per 1,000 vs Raw Savings Per Beneficiary<br>Bubble size = Overall Weighted Risk Score (PY)",
                labels={
                    "P_EM_PCP_Vis": "PCP E&M Visits per 1,000 Person-Years",
                    "Raw_Savings_Per_Benef": "Raw Savings Per Beneficiary ($)"
                }
            )
            # Add callout annotation with actual ACO name
            selected_row = plot_df[plot_df["ACO_ID"] == aco_data["ACO_ID"]]
            if not selected_row.empty:
                x_val = selected_row["P_EM_PCP_Vis"].iloc[0]
                y_val = selected_row["Raw_Savings_Per_Benef"].iloc[0]
                aco_name = selected_row["ACO_Name"].iloc[0]
                fig_scat_per.add_annotation(
                    x=x_val,
                    y=y_val,
                    text=f"<b>{aco_name}</b>",
                    showarrow=True,
                    arrowhead=2,
                    ax=60,
                    ay=-60,
                    font=dict(size=14, color="#000000"),
                    bgcolor="#FFFFFF",
                    bordercolor="#000000",
                    borderwidth=2,
                    arrowwidth=2,
                    arrowcolor="#000000"
                )
            fig_scat_per.update_traces(marker=dict(opacity=0.7, line=dict(width=0.5)))
            fig_scat_per.update_layout(template=PLOTLY_TEMPLATE, height=600)
            st.plotly_chart(fig_scat_per, use_container_width=True)
        else:
            st.info("Missing columns for per-beneficiary scatter plot (P_EM_PCP_Vis, BnchmkMinExp, N_AB_Year_PY, ACO_ID).")
        # 6. Utilization Comparison – FIXED ORDER VERSION
        st.markdown("### Utilization Comparison")
        st.write("Bar chart of per capita cost across all ACOs (ascending order: lowest/left = most efficient → highest/right = least efficient). Selected ACO highlighted in theme color.")
        for cat, cost_col in capann_cols.items():
            if cost_col not in df.columns:
                continue
            # Copy and prepare
            all_df = df.copy()
            all_df[cost_col] = pd.to_numeric(all_df[cost_col], errors='coerce')
            # Sort ascending (lowest cost first = left)
            all_df = all_df.sort_values(cost_col, ascending=True, na_position='last').reset_index(drop=True)
            # Highlight/color AFTER sort
            all_df["Highlight"] = all_df["ACO_ID"] == aco_data["ACO_ID"]
            all_df["Color"] = all_df["Highlight"].map({True: PRIMARY, False: NEUTRAL}) # ← use PRIMARY for highlight, NEUTRAL for others
            # Add index as x (forces correct sorted order)
            all_df["Index"] = all_df.index
            fig_bar = px.bar(
                all_df,
                x="Index", # Use numeric index instead of name
                y=cost_col,
                color="Color",
                color_discrete_map={PRIMARY: PRIMARY, NEUTRAL: NEUTRAL}, # ← map to theme colors
                title=f"{cat} – Per Capita Cost Across All ACOs (Selected ACO in Theme Color)",
                labels={cost_col: "Per Capita Cost ($)"},
                hover_data=["ACO_Name", cost_col]
            )
            # Custom x-axis: only show selected ACO name at its position
            ticktext = ["" for _ in all_df.index] # blank everywhere
            if all_df["Highlight"].any():
                selected_idx = all_df[all_df["Highlight"]].index[0]
                ticktext[selected_idx] = aco_data["ACO_Name"]
            fig_bar.update_layout(
                template=PLOTLY_TEMPLATE,
                xaxis={
                    'tickmode': 'array',
                    'tickvals': list(range(len(all_df))), # 0,1,2,... positions
                    'ticktext': ticktext,
                },
                xaxis_tickangle=-45,
                showlegend=False,
                height=500
            )
            # Fixed hover: ACO name + formatted cost
            fig_bar.update_traces(hovertemplate="%{customdata[0]}<br>Per Capita Cost: $%{y:,.0f}")
            st.plotly_chart(fig_bar, use_container_width=True)
        # New: 4. Outlier Detection in Admissions/LOS
        st.markdown("### Outlier Detection in Admissions & LOS")
        st.write("Services where the ACO deviates >20% from track average (threshold for highlighting).")
        outlier_data = []
        outlier_threshold = 0.20 # 20% deviation
        st.caption(f"**Threshold**: >{outlier_threshold*100}% deviation from track average (absolute value).")
        # Inpatient Admissions per 1,000
        if all(col in aco_data for col in ["ADM_INP_All", "N_AB"]) and aco_data["N_AB"] > 0:
            adm_per_1000 = aco_data["ADM_INP_All"] / aco_data["N_AB"] * 1000
            track_adm_per_1000 = track_avg["ADM_INP_All"] / track_avg["N_AB"] * 1000 if track_avg["N_AB"] > 0 else np.nan
            pct_diff = (adm_per_1000 - track_adm_per_1000) / track_adm_per_1000 if pd.notna(track_adm_per_1000) and track_adm_per_1000 != 0 else np.nan
            if abs(pct_diff) > outlier_threshold:
                outlier_data.append({
                    "Service": "Inpatient Admissions per 1,000",
                    "This ACO": f"{adm_per_1000:.1f}",
                    "Track Avg": f"{track_adm_per_1000:.1f}",
                    "% Deviation": f"{pct_diff*100:.1f}%",
                    "Alert": "High" if pct_diff > 0 else "Low"
                })
        # Outpatient Visits per 1,000
        if all(col in aco_data for col in ["ADM_OPD", "N_AB"]) and aco_data["N_AB"] > 0:
            adm_per_1000 = aco_data["ADM_OPD"] / aco_data["N_AB"] * 1000
            track_adm_per_1000 = track_avg["ADM_OPD"] / track_avg["N_AB"] * 1000 if track_avg["N_AB"] > 0 else np.nan
            pct_diff = (adm_per_1000 - track_adm_per_1000) / track_adm_per_1000 if pd.notna(track_adm_per_1000) and track_adm_per_1000 != 0 else np.nan
            if abs(pct_diff) > outlier_threshold:
                outlier_data.append({
                    "Service": "Outpatient Visits per 1,000",
                    "This ACO": f"{adm_per_1000:.1f}",
                    "Track Avg": f"{track_adm_per_1000:.1f}",
                    "% Deviation": f"{pct_diff*100:.1f}%",
                    "Alert": "High" if pct_diff > 0 else "Low"
                })
        # Physician/Supplier Visits per 1,000
        if all(col in aco_data for col in ["ADM_PB", "N_AB"]) and aco_data["N_AB"] > 0:
            adm_per_1000 = aco_data["ADM_PB"] / aco_data["N_AB"] * 1000
            track_adm_per_1000 = track_avg["ADM_PB"] / track_avg["N_AB"] * 1000 if track_avg["N_AB"] > 0 else np.nan
            pct_diff = (adm_per_1000 - track_adm_per_1000) / track_adm_per_1000 if pd.notna(track_adm_per_1000) and track_adm_per_1000 != 0 else np.nan
            if abs(pct_diff) > outlier_threshold:
                outlier_data.append({
                    "Service": "Physician/Supplier Visits per 1,000",
                    "This ACO": f"{adm_per_1000:.1f}",
                    "Track Avg": f"{track_adm_per_1000:.1f}",
                    "% Deviation": f"{pct_diff*100:.1f}%",
                    "Alert": "High" if pct_diff > 0 else "Low"
                })
        # SNF LOS
        if "SNF_LOS" in aco_data and pd.notna(aco_data["SNF_LOS"]):
            snf_los = aco_data["SNF_LOS"]
            track_snf_los = track_avg["SNF_LOS"] if pd.notna(track_avg["SNF_LOS"]) else np.nan
            pct_diff = (snf_los - track_snf_los) / track_snf_los if pd.notna(track_snf_los) and track_snf_los != 0 else np.nan
            if abs(pct_diff) > outlier_threshold:
                outlier_data.append({
                    "Service": "SNF Length of Stay (days)",
                    "This ACO": f"{snf_los:.1f}",
                    "Track Avg": f"{track_snf_los:.1f}",
                    "% Deviation": f"{pct_diff*100:.1f}%",
                    "Alert": "High" if pct_diff > 0 else "Low"
                })
        if outlier_data:
            outlier_df = pd.DataFrame(outlier_data)
            st.dataframe(outlier_df, use_container_width=True, hide_index=True)
        else:
            st.success("No significant outliers detected (>20% deviation from track average).")
        # New: 6. Savings Attribution to Utilization Reduction
        st.markdown("### Savings Attribution to Utilization Reduction")
        st.write("Estimated contribution of lower utilization to savings (compared to track average). Calculation: (Track Avg Per Capita - ACO Per Capita) × N_AB. Positive = contributed to savings; negative = increased costs.")
        attribution_data = []
        for cat, cost_col in capann_cols.items():
            if cost_col not in df.columns or pd.isna(aco_data[cost_col]) or pd.isna(track_avg[cost_col]):
                continue
            delta_per_capita = track_avg[cost_col] - aco_data[cost_col]
            contribution = delta_per_capita * aco_data["N_AB"] if aco_data["N_AB"] > 0 else 0
            attribution_data.append({
                "Service Category": cat,
                "Delta vs Track Avg (Per Capita)": fmt_dollars(delta_per_capita, 0),
                "Estimated Contribution to Savings": fmt_dollars(contribution, 0),
                "Note": "Positive = lower utilization helped savings" if contribution > 0 else "Negative = higher utilization reduced savings"
            })
        if attribution_data:
            attribution_df = pd.DataFrame(attribution_data)
            st.dataframe(attribution_df, use_container_width=True, hide_index=True)
            st.caption("**Calculation Note**: Contribution = (Track Avg Per Capita - ACO Per Capita) × Assigned Beneficiaries (N_AB). Assumes utilization delta directly impacts expenditures; actual savings also depend on benchmark and quality performance.")
        else:
            st.info("No valid utilization data for attribution analysis.")
        # 1. Quality Gates & Flags
        st.markdown("### Quality Gates & Program Flags")
        st.write("Key indicators for meeting MSSP quality requirements and adjustments (PY 2023).")
        quality_gates = [
            {"Flag": "Extreme & Uncontrollable Circumstance (Quality)", "Value": "Yes" if aco_data.get("DisAffQual", 0) == 1 else "No"},
            {"Flag": "Met Quality Performance Standard", "Value": "Yes" if aco_data.get("Met_QPS", 0) == 1 else "No"},
            {"Flag": "Met Alternative Quality Performance Standard", "Value": "Yes" if aco_data.get("Met_AltQPS", 0) == 1 else "No"},
            {"Flag": "Met or Exceeded 30th Percentile MIPS QPC Score", "Value": "Yes" if aco_data.get("Met_30pctl", 0) == 1 else "No"},
            {"Flag": "Met eCQM/MIPS CQM Reporting Incentive", "Value": "Yes" if aco_data.get("Met_Incentive", 0) == 1 else "No"},
            {"Flag": "1st Year ACO Met Reporting Criteria", "Value": "Yes" if aco_data.get("Met_FirstYear", 0) == 1 else "No"},
            {"Flag": "Reported CMS Web Interface Measure Set", "Value": "Yes" if aco_data.get("Report_WI", 0) == 1 else "No"},
            {"Flag": "Reported eCQMs or MIPS CQMs", "Value": "Yes" if aco_data.get("Report_eCQM_CQM", 0) == 1 else "No"},
            {"Flag": "Incomplete Reporting", "Value": "Yes" if aco_data.get("Report_Inc", 0) == 1 else "No"},
            {"Flag": "Extreme & Uncontrollable Circumstance - 30th Percentile Adjustment", "Value": "Yes" if aco_data.get("Recvd30p", 0) == 1 else "No"},
        ]
        df_gates = pd.DataFrame(quality_gates)
        st.dataframe(df_gates, use_container_width=True, hide_index=True)
        # 2. Quality Program Measures
        st.markdown("### Quality Program Measures")
        st.write("Individual quality measures with scores, track average, all-ACOs average, and % difference. Readmission/admission rates formatted as % (lower is better).")
        quality_measures = [
            {"Measure": "Hospital-wide 30-day Readmission Rate", "col": "Measure_479", "lower_better": True, "is_rate": True},
            {"Measure": "All-cause Unplanned Admissions for Patients with Multiple Chronic Conditions", "col": "Measure_484", "lower_better": True, "is_rate": False},
            {"Measure": "Falls: Screening for Future Fall Risk", "col": "QualityID_318", "lower_better": False, "is_rate": False},
            {"Measure": "Preventive Care & Screening: Influenza Immunization", "col": "QualityID_110", "lower_better": False, "is_rate": False},
            {"Measure": "Preventive Care & Screening: Tobacco Use Screening & Cessation Intervention", "col": "QualityID_226", "lower_better": False, "is_rate": False},
            {"Measure": "Preventive Care & Screening: Depression Screening & Follow-up Plan (WI)", "col": "QualityID_134_WI", "lower_better": False, "is_rate": False},
            {"Measure": "Colorectal Cancer Screening", "col": "QualityID_113", "lower_better": False, "is_rate": False},
            {"Measure": "Breast Cancer Screening", "col": "QualityID_112", "lower_better": False, "is_rate": False},
            {"Measure": "Statin Therapy for Prevention & Treatment of Cardiovascular Disease", "col": "QualityID_438", "lower_better": False, "is_rate": False},
            {"Measure": "Depression Remission at Twelve Months", "col": "QualityID_370", "lower_better": False, "is_rate": False},
            {"Measure": "Diabetes: Hemoglobin A1c Poor Control (>9%)", "col": "QualityID_001_WI", "lower_better": True, "is_rate": False},
            {"Measure": "Controlling High Blood Pressure", "col": "QualityID_236_WI", "lower_better": False, "is_rate": False},
        ]
        quality_table = []
        chart_data = [] # For bar chart
        for m in quality_measures:
            col = m["col"]
            if col not in df.columns:
                continue
            this_aco_val = pd.to_numeric(aco_data.get(col, np.nan), errors='coerce')
            track_avg_val = track_avg.get(col, np.nan)
            all_aco_avg_val = df[col].mean(skipna=True)
            pct_diff_track = (this_aco_val - track_avg_val) / track_avg_val if pd.notna(track_avg_val) and track_avg_val != 0 else np.nan
            pct_diff_all = (this_aco_val - all_aco_avg_val) / all_aco_avg_val if pd.notna(all_aco_avg_val) and all_aco_avg_val != 0 else np.nan
            # Format as % for rate measures, otherwise as %
            if m["is_rate"]:
                this_fmt = f"{this_aco_val*100:.2f}%" if pd.notna(this_aco_val) else "N/A"
                track_fmt = f"{track_avg_val*100:.2f}%" if pd.notna(track_avg_val) else "N/A"
                all_fmt = f"{all_aco_avg_val*100:.2f}%" if pd.notna(all_aco_avg_val) else "N/A"
            else:
                this_fmt = fmt_pct(this_aco_val) if pd.notna(this_aco_val) else "N/A"
                track_fmt = fmt_pct(track_avg_val) if pd.notna(track_avg_val) else "N/A"
                all_fmt = fmt_pct(all_aco_avg_val) if pd.notna(all_aco_avg_val) else "N/A"
            quality_table.append({
                "Measure": m["Measure"],
                "This ACO": this_fmt,
                "Track Avg": track_fmt,
                "% Diff vs Track": f"{pct_diff_track*100:.1f}%" if pd.notna(pct_diff_track) else "N/A",
                "All ACOs Avg": all_fmt,
                "% Diff vs All": f"{pct_diff_all*100:.1f}%" if pd.notna(pct_diff_all) else "N/A"
            })
            # Add to chart data (use raw values for plotting)
            chart_data.append({
                "Measure": m["Measure"],
                "This ACO": this_aco_val,
                "Track Avg": track_avg_val,
                "All ACOs Avg": all_aco_avg_val,
                "Lower Better": m["lower_better"]
            })
        if quality_table:
            df_quality = pd.DataFrame(quality_table)
            st.dataframe(df_quality, use_container_width=True, hide_index=True)
            # Bar chart visualization
            if chart_data:
                df_chart = pd.DataFrame(chart_data)
                df_chart_melt = df_chart.melt(id_vars=["Measure", "Lower Better"], var_name="Comparison", value_name="Score")
                fig_quality = px.bar(
                    df_chart_melt,
                    x="Measure",
                    y="Score",
                    color="Comparison",
                    barmode="group",
                    title="Quality Measures Comparison",
                    labels={"Score": "Score (%) or Rate"},
                    height=500
                )
                fig_quality.update_layout(template=PLOTLY_TEMPLATE, xaxis_tickangle=-45)
                st.plotly_chart(fig_quality, use_container_width=True)
        else:
            st.info("No quality measure columns found.")
        # 3. Patient Experience (CAHPS) Measures
        st.markdown("### Patient Experience (CAHPS) Measures")
        st.write("CAHPS domain scores (higher is better) vs track and all-ACOs averages.")
        cahps_map = {
            "CAHPS_1": "Getting Timely Care, Appointments, and Information",
            "CAHPS_2": "How Well Providers Communicate",
            "CAHPS_3": "Patients' Rating of Provider",
            "CAHPS_4": "Access to Specialists",
            "CAHPS_5": "Health Promotion and Education",
            "CAHPS_6": "Shared Decision Making",
            "CAHPS_7": "Health Status/Functional Status",
            "CAHPS_11": "Stewardship of Patient Resources",
            "CAHPS_9": "Courteous and Helpful Office Staff",
            "CAHPS_8": "Care Coordination"
        }
        cahps_table = []
        chart_data = [] # For bar chart
        for raw_col, friendly_name in cahps_map.items():
            col = raw_col
            if col not in df.columns:
                continue
            this_aco_val = pd.to_numeric(aco_data.get(col, np.nan), errors='coerce')
            track_avg_val = track_avg.get(col, np.nan)
            all_aco_avg_val = df[col].mean(skipna=True)
            pct_diff_track = (this_aco_val - track_avg_val) / track_avg_val if pd.notna(track_avg_val) and track_avg_val != 0 else np.nan
            pct_diff_all = (this_aco_val - all_aco_avg_val) / all_aco_avg_val if pd.notna(all_aco_avg_val) and all_aco_avg_val != 0 else np.nan
            cahps_table.append({
                "CAHPS Domain": friendly_name,
                "This ACO": fmt_pct(this_aco_val) if pd.notna(this_aco_val) else "N/A",
                "Track Avg": fmt_pct(track_avg_val) if pd.notna(track_avg_val) else "N/A",
                "% Diff vs Track": f"{pct_diff_track*100:.1f}%" if pd.notna(pct_diff_track) else "N/A",
                "All ACOs Avg": fmt_pct(all_aco_avg_val) if pd.notna(all_aco_avg_val) else "N/A",
                "% Diff vs All": f"{pct_diff_all*100:.1f}%" if pd.notna(pct_diff_all) else "N/A"
            })
            chart_data.append({
                "CAHPS Domain": friendly_name,
                "This ACO": this_aco_val,
                "Track Avg": track_avg_val,
                "All ACOs Avg": all_aco_avg_val
            })
        if cahps_table:
            df_cahps = pd.DataFrame(cahps_table)
            st.dataframe(df_cahps, use_container_width=True, hide_index=True)
            # Bar chart visualization
            if chart_data:
                df_chart = pd.DataFrame(chart_data)
                df_chart_melt = df_chart.melt(id_vars="CAHPS Domain", var_name="Comparison", value_name="Score")
                fig_cahps = px.bar(
                    df_chart_melt,
                    x="CAHPS Domain",
                    y="Score",
                    color="Comparison",
                    barmode="group",
                    title="CAHPS Domain Scores Comparison",
                    labels={"Score": "Score (%)"},
                    height=500
                )
                fig_cahps.update_layout(template=PLOTLY_TEMPLATE, xaxis_tickangle=-45, yaxis_range=[0, 100])
                st.plotly_chart(fig_cahps, use_container_width=True)
        else:
            st.info("No CAHPS columns found in data.")

elif tab_selection == "Program Changes":
    st.subheader("MSSP Program Changes: PY 2022 to PY 2023")

    st.dataframe(df_changes, use_container_width=True, hide_index=True)

